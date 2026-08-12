"""Build the project deck as a .pptx.

Run:  .venv/bin/python scripts/build_deck.py [--shots DIR] [--out FILE]

The deck is generated rather than hand-made so it can be rebuilt after the UI
changes: recapture the screenshots, re-run this, and the slides stay in sync.
Fonts are deliberately Arial and Courier New — they exist on every machine that
might open the file, which matters more here than typographic preference.
"""

from __future__ import annotations

import argparse
import math
from pathlib import Path

try:
    from pptx import Presentation
except ModuleNotFoundError as exc:  # pragma: no cover - a tooling dependency
    raise SystemExit(
        "python-pptx is not installed. Run:  .venv/bin/pip install python-pptx"
    ) from exc

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.table import _Cell
from pptx.text.text import TextFrame
from pptx.util import Emu, Inches, Length, Pt

Stat = tuple[str, str, RGBColor]
Item = str | tuple[str, str]
#: endpoint, name, class, method, output, model-cost, colour, is-called
ArchStage = tuple[str, str, str, str, str, str, RGBColor, bool]

# --- Palette, taken from the application's own tokens ------------------------

BG = RGBColor(0x0A, 0x0F, 0x1C)
SURFACE = RGBColor(0x11, 0x18, 0x27)
RULE = RGBColor(0x1F, 0x2A, 0x40)
TEXT = RGBColor(0xE8, 0xED, 0xF5)
DIM = RGBColor(0x88, 0x96, 0xB0)
MUTED = RGBColor(0x4B, 0x59, 0x75)
ACCENT = RGBColor(0x60, 0xA5, 0xFA)
GREEN = RGBColor(0x34, 0xD3, 0x99)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)

SANS = "Arial"
MONO = "Courier New"

W = Inches(13.333)
H = Inches(7.5)
MARGIN = Inches(0.75)


# --- Primitives --------------------------------------------------------------


def blank(prs: Presentation) -> Slide:
    """A slide with no placeholders, painted with the deck background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, W, H)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def textbox(slide: Slide, left: Length, top: Length, width: Length, height: Length) -> TextFrame:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def para(
    frame: TextFrame,
    text: str,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
    font: str = SANS,
    space_after: int = 6,
    first: bool = False,
) -> None:
    p = frame.paragraphs[0] if first else frame.add_paragraph()
    p.text = text
    p.space_after = Pt(space_after)
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font


def eyebrow(slide: Slide, text: str, top: Length | None = None) -> None:
    frame = textbox(slide, MARGIN, top or Inches(0.5), Inches(11), Inches(0.3))
    para(frame, text.upper(), size=11, color=MUTED, bold=True, font=MONO, first=True)


def heading(slide: Slide, text: str, top: Length | None = None, size: int = 30) -> None:
    frame = textbox(slide, MARGIN, top or Inches(0.85), Inches(11.8), Inches(1.0))
    para(frame, text, size=size, color=TEXT, bold=True, first=True)


def rule(slide: Slide, top: Length, width: Length | None = None, color: RGBColor = RULE) -> None:
    line = slide.shapes.add_shape(1, MARGIN, top, width or Inches(11.8), Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False


def _cell_borders(cell: _Cell, colour: RGBColor = RULE, width_pt: float = 1.0) -> None:
    """Set a cell's four borders.

    python-pptx has no API for this, and the default a renderer picks is a white
    hairline — which is loud on a dark deck. Written straight into the cell's
    XML properties instead.
    """
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for existing in tc_pr.findall(qn(tag)):
            tc_pr.remove(existing)
        ln = OxmlElement(tag)
        ln.set("w", str(int(width_pt * 12700)))
        ln.set("cap", "flat")
        fill = OxmlElement("a:solidFill")
        clr = OxmlElement("a:srgbClr")
        clr.set("val", f"{colour}")
        fill.append(clr)
        ln.append(fill)
        tc_pr.append(ln)


def card(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    *,
    accent: RGBColor | None = None,
) -> None:
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = SURFACE
    box.line.color.rgb = accent or RULE
    box.line.width = Pt(1.25)
    box.shadow.inherit = False


# --- Slide kinds -------------------------------------------------------------


def title_slide(prs: Presentation, title: str, subtitle: str, footer: str) -> None:
    slide = blank(prs)
    bar = slide.shapes.add_shape(1, 0, Inches(2.62), Inches(0.09), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    frame = textbox(slide, MARGIN, Inches(2.5), Inches(11.5), Inches(2.4))
    para(frame, title, size=46, color=TEXT, bold=True, space_after=14, first=True)
    para(frame, subtitle, size=19, color=DIM, space_after=0)

    foot = textbox(slide, MARGIN, Inches(6.5), Inches(11.5), Inches(0.5))
    para(foot, footer, size=12, color=MUTED, font=MONO, first=True)


def bullets_slide(
    prs: Presentation,
    eyebrow_text: str,
    title: str,
    items: list[Item],
    *,
    note: str | None = None,
) -> None:
    slide = blank(prs)
    eyebrow(slide, eyebrow_text)
    heading(slide, title)
    rule(slide, Inches(1.85))

    frame = textbox(slide, MARGIN, Inches(2.15), Inches(11.8), Inches(4.4))
    first = True
    for item in items:
        if isinstance(item, tuple):
            lead, body = item
            para(frame, lead, size=17, color=TEXT, bold=True, space_after=3, first=first)
            para(frame, body, size=14.5, color=DIM, space_after=16)
        else:
            para(frame, item, size=16, color=DIM, space_after=12, first=first)
        first = False

    if note:
        nf = textbox(slide, MARGIN, Inches(6.55), Inches(11.8), Inches(0.6))
        para(nf, note, size=12.5, color=DIM, font=MONO, first=True)


def image_slide(
    prs: Presentation, eyebrow_text: str, title: str, image_path: Path, caption: str
) -> None:
    """Full-width screenshot with a title above and one line of context below."""
    slide = blank(prs)
    eyebrow(slide, eyebrow_text, top=Inches(0.42))
    frame = textbox(slide, MARGIN, Inches(0.72), Inches(11.8), Inches(0.5))
    para(frame, title, size=24, color=TEXT, bold=True, first=True)

    pic_h = Inches(5.22)
    pic_w = Emu(int(pic_h * 16 / 9))
    left = Emu(int((W - pic_w) / 2))
    slide.shapes.add_picture(str(image_path), left, Inches(1.3), width=pic_w, height=pic_h)

    cap = textbox(slide, MARGIN, Inches(6.68), Inches(11.8), Inches(0.55))
    para(cap, caption, size=13, color=DIM, first=True)


def stat_slide(
    prs: Presentation,
    eyebrow_text: str,
    title: str,
    stats: list[Stat],
    footer: str | None = None,
) -> None:
    slide = blank(prs)
    eyebrow(slide, eyebrow_text)
    heading(slide, title)
    rule(slide, Inches(1.85))

    cols = len(stats)
    gap = Inches(0.3)
    total = Inches(11.8)
    width = Emu(int((total - gap * (cols - 1)) / cols))
    for i, (value, label, colour) in enumerate(stats):
        left = Emu(int(MARGIN + i * (width + gap)))
        card(slide, left, Inches(2.35), width, Inches(2.0))
        vf = textbox(
            slide,
            Emu(int(left + Inches(0.28))),
            Inches(2.7),
            Emu(int(width - Inches(0.56))),
            Inches(1.0),
        )
        para(vf, value, size=40, color=colour, bold=True, first=True)
        lf = textbox(
            slide,
            Emu(int(left + Inches(0.28))),
            Inches(3.55),
            Emu(int(width - Inches(0.56))),
            Inches(0.7),
        )
        para(lf, label, size=12.5, color=DIM, first=True)

    if footer:
        ff = textbox(slide, MARGIN, Inches(5.0), Inches(11.8), Inches(1.2))
        para(ff, footer, size=15, color=DIM, first=True)


def _arrowhead(
    shape: BaseShape, colour: RGBColor, width_pt: float = 1.5, dash: bool = False
) -> None:
    """Stroke a line and put a triangular head on its end.

    python-pptx exposes line colour and width but not arrowheads, so the
    `a:tailEnd` element is written into the line properties directly.
    """
    line = shape.line
    line.color.rgb = colour
    line.width = Pt(width_pt)
    ln = line._get_or_add_ln()
    if dash:
        dash_el = OxmlElement("a:prstDash")
        dash_el.set("val", "dash")
        ln.append(dash_el)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    ln.append(tail)


def arrow(slide: Slide, x1: Length, y1: Length, x2: Length, y2: Length, colour: RGBColor) -> None:
    """A straight connector with an arrowhead at the far end."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    _arrowhead(conn, colour)


def dotted_drop(slide: Slide, x: Length, y1: Length, y2: Length) -> None:
    """A faint vertical tie between a stage and the band below it."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    conn.line.color.rgb = RULE
    conn.line.width = Pt(1.0)
    dash = OxmlElement("a:prstDash")
    dash.set("val", "sysDot")
    conn.line._get_or_add_ln().append(dash)


def loop_arrow(slide: Slide, left: Length, right: Length, top: Length, rise: Length) -> None:
    """The self-loop drawn over a stage that repeats.

    A semicircle sampled into line segments. PowerPoint's arc primitive takes no
    arrowhead without more XML than the shape is worth, and a four-segment
    approximation reads as a rendering glitch rather than a loop.
    """
    steps = 24
    centre = (left + right) / 2
    radius_x = (right - left) * 0.26
    points = []
    for step in range(1, steps + 1):
        angle = math.pi * step / steps
        points.append(
            (
                Emu(int(centre + radius_x * math.cos(angle))),
                Emu(int(top - rise * math.sin(angle))),
            )
        )
    builder = slide.shapes.build_freeform(Emu(int(centre + radius_x)), top)
    builder.add_line_segments(points, close=False)
    shape = builder.convert_to_shape()
    shape.fill.background()
    _arrowhead(shape, DIM, width_pt=1.25)


def pill(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    text: str,
    colour: RGBColor,
    *,
    dashed: bool = False,
) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.18
    box.fill.background()
    box.line.color.rgb = colour
    box.line.width = Pt(1.5 if not dashed else 1.0)
    box.shadow.inherit = False
    if dashed:
        dash = OxmlElement("a:prstDash")
        dash.set("val", "dash")
        box.line._get_or_add_ln().append(dash)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    para(frame, text, size=12, color=colour, bold=not dashed, first=True)
    frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def architecture_slide(
    prs: Presentation, eyebrow_text: str, title: str, stages: list[ArchStage], footer: str
) -> None:
    """The pipeline drawn end to end: endpoint, stage, output, and model cost."""
    slide = blank(prs)
    eyebrow(slide, eyebrow_text)
    heading(slide, title, size=28)

    n = len(stages)
    gap = Inches(0.2)
    total = Inches(11.8)
    width = Emu(int((total - gap * (n - 1)) / n))
    box_top, box_h = Inches(2.46), Inches(1.28)
    mid = Emu(int(box_top + box_h / 2))

    for i, (endpoint, name, cls, method, output, cost, colour, always) in enumerate(stages):
        left = Emu(int(MARGIN + i * (width + gap)))
        centre = Emu(int(left + width / 2))
        inner = Emu(int(width - Inches(0.24)))
        pad = Emu(int(left + Inches(0.12)))

        ef = textbox(slide, pad, Inches(2.02), inner, Inches(0.3))
        para(ef, endpoint, size=10.5, color=MUTED, font=MONO, first=True)

        card(slide, left, box_top, width, box_h)
        nf = textbox(slide, pad, Inches(2.60), inner, Inches(0.34))
        para(nf, name, size=17, color=TEXT, bold=True, first=True)
        nf.paragraphs[0].alignment = PP_ALIGN.CENTER
        cf = textbox(slide, pad, Inches(3.00), inner, Inches(0.6))
        para(cf, cls, size=10.5, color=DIM, font=MONO, first=True)
        cf.paragraphs[0].alignment = PP_ALIGN.CENTER
        para(cf, method, size=10.5, color=DIM, font=MONO, space_after=0)
        cf.paragraphs[1].alignment = PP_ALIGN.CENTER

        of = textbox(slide, pad, Inches(3.88), inner, Inches(0.32))
        para(of, output, size=10.5, color=MUTED, font=MONO, first=True)
        of.paragraphs[0].alignment = PP_ALIGN.CENTER

        if i < n - 1:
            arrow(
                slide,
                Emu(int(left + width + Inches(0.02))),
                mid,
                Emu(int(left + width + gap - Inches(0.02))),
                mid,
                DIM,
            )

        dotted_drop(slide, centre, Inches(4.26), Inches(5.00))
        pill(
            slide,
            left,
            Inches(5.00),
            width,
            Inches(0.62),
            cost,
            colour if always else MUTED,
            dashed=not always,
        )

    # Refine repeats until the draft has no open questions.
    refine_left = Emu(int(MARGIN + (width + gap)))
    loop_arrow(slide, refine_left, Emu(int(refine_left + width)), box_top, Inches(0.30))
    lf = textbox(slide, Emu(int(refine_left - Inches(0.6))), Inches(1.58), Inches(3.3), Inches(0.3))
    para(lf, "until no questions remain", size=10.5, color=DIM, first=True)
    lf.paragraphs[0].alignment = PP_ALIGN.CENTER

    rule(slide, Inches(4.64))
    bf = textbox(slide, MARGIN, Inches(4.71), Inches(6.0), Inches(0.3))
    para(bf, "WHERE THE MODEL IS CALLED", size=10.5, color=MUTED, font=MONO, first=True)

    ff = textbox(slide, MARGIN, Inches(5.94), Inches(11.8), Inches(1.0))
    para(ff, footer, size=13.5, color=DIM, first=True)


def table_slide(
    prs: Presentation,
    eyebrow_text: str,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    standfirst: str | None = None,
    note: str | None = None,
) -> None:
    slide = blank(prs)
    eyebrow(slide, eyebrow_text)
    heading(slide, title, size=27)

    top = Inches(2.2)
    if standfirst:
        sf = textbox(slide, MARGIN, Inches(1.62), Inches(11.2), Inches(0.8))
        para(sf, standfirst, size=15, color=DIM, first=True)
        top = Inches(2.46)
    else:
        rule(slide, Inches(1.85))

    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), MARGIN, top, Inches(11.8), Inches(0.4)
    )
    table = shape.table
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _cell_borders(cell)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = DIM
                run.font.name = MONO
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            _cell_borders(cell)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(13)
                    run.font.color.rgb = TEXT if c == 0 else DIM
                    run.font.name = SANS
                    run.font.bold = c == 0

    if note:
        nf = textbox(slide, MARGIN, Inches(6.5), Inches(11.8), Inches(0.6))
        para(nf, note, size=12.5, color=DIM, font=MONO, first=True)


def closing_slide(prs: Presentation, title: str, lines: list[str], footer: str) -> None:
    slide = blank(prs)
    frame = textbox(slide, MARGIN, Inches(2.3), Inches(11.5), Inches(1.0))
    para(frame, title, size=38, color=TEXT, bold=True, first=True)
    body = textbox(slide, MARGIN, Inches(3.4), Inches(10.5), Inches(2.6))
    first = True
    for line in lines:
        para(body, line, size=17, color=DIM, space_after=14, first=first)
        first = False
    foot = textbox(slide, MARGIN, Inches(6.5), Inches(11.5), Inches(0.5))
    para(foot, footer, size=12, color=MUTED, font=MONO, first=True)


# --- The deck ----------------------------------------------------------------


def build(shots: Path, out: Path) -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    title_slide(
        prs,
        "Data Source Onboarding Agent",
        "Plain English in. A validated, standards-compliant, versioned Python connector out.",
        "natural language → spec → generated code → machine-checked → proven against "
        "a live database",
    )

    # --- 2. Business case, in one slide --------------------------------------

    table_slide(
        prs,
        "Business objective",
        "The repetition is where the risk lives",
        ["Step in onboarding", "Today", "With the agent"],
        [
            [
                "Capture the requirement",
                "A ticket, then a clarifying thread over days",
                "Plain English; the gaps are asked about in the same session",
            ],
            [
                "Write the connector",
                "Hand-written per source, per engineer",
                "Rendered from a reviewed spec — identical every time",
            ],
            [
                "Enforce the standards",
                "Code review, after the code already exists",
                "13 machine checks, before it can be accepted at all",
            ],
            [
                "Prove it connects",
                "Run it locally and hope the target matches",
                "Sandboxed test against the real database, recorded",
            ],
            [
                "Answer “why is this in production?”",
                "Ask whoever wrote it, if they still work here",
                "A manifest chaining request → spec → template → checksum",
            ],
        ],
        standfirst="The cost is not writing one connector — it is writing the fortieth, where the "
        "same five methods get re-derived slightly differently and the standards hold only as "
        "well as the reviewer's attention that afternoon.",
        note="Requesting team: self-service · Engineering: reviews a diff, not a blank file · "
        "Security: credentials structurally excluded · Audit: every artifact traceable",
    )

    # --- 3. The guarantee ----------------------------------------------------

    bullets_slide(
        prs,
        "The core idea",
        "Reliability comes from constraint, not from prompting",
        [
            (
                "The model never writes the code",
                "It fills a structured spec — schema-enforced, so it selects from closed sets "
                "rather than inventing strings. A Jinja2 template renders the connector, "
                "deterministically: same spec in, same bytes out.",
            ),
            (
                "13 machine checks decide what ships",
                "no-hardcoded-credentials · no-dynamic-sql · env-for-secrets · "
                "no-dangerous-calls · type-hints · docstrings, and seven more — read from the "
                "parsed syntax tree, never "
                "by importing the code, because importing runs it.",
            ),
            (
                "Repair is bounded and cannot regress",
                "At most three attempts, and a candidate is adopted only if it strictly "
                "reduces the error count. A rejected connector is still returned, with its "
                "findings, so a "
                "reviewer can see exactly what went wrong.",
            ),
            (
                "Secrets cannot reach the artifact",
                "Credentials are scrubbed from the request before the model sees it, two checks "
                "make a literal credential a rejection, and the manifest records the environment "
                "variables by name only — never their values.",
            ),
        ],
        note="A checker that has only ever seen good input has demonstrated nothing — so the suite "
        "injects faults on purpose.",
    )

    # --- 4. Architecture -----------------------------------------------------

    architecture_slide(
        prs,
        "Architecture",
        "One request, five stages",
        [
            (
                "POST /api/requests",
                "Extract",
                "SpecExtractor",
                ".extract()",
                "SpecDraft + questions",
                "always — 1 call",
                ACCENT,
                True,
            ),
            (
                "POST …/answers",
                "Refine",
                "SpecExtractor",
                ".refine()",
                "SpecDraft (complete)",
                "never",
                GREEN,
                False,
            ),
            (
                "POST …/generate",
                "Generate",
                "ConnectorGenerator",
                ".generate()",
                "GeneratedConnector",
                "only on failure / docs",
                ACCENT,
                True,
            ),
            (
                "POST …/test",
                "Test",
                "ConnectionSandbox",
                ".run()",
                "SandboxResult",
                "only on failure",
                ACCENT,
                True,
            ),
            (
                "GET …/download",
                "Deliver",
                "Artifact",
                ".to_zip()",
                "bytes",
                "never",
                GREEN,
                False,
            ),
        ],
        footer="Answering a clarifying question is deterministic — the answers map onto known "
        "fields, so no second extraction can misread what you already settled.",
    )

    # --- 5. Technology -------------------------------------------------------

    table_slide(
        prs,
        "Technology",
        "The stack, and why each piece is there",
        ["Layer", "Choice", "Why this one"],
        [
            [
                "API",
                "FastAPI + Pydantic v2",
                "Schema validation at the boundary, not in the handlers",
            ],
            ["Spec", "Pydantic models", "Permissive draft, then promotion to a strict spec"],
            ["Codegen", "Jinja2", "Deterministic. The reason the standards guarantee holds"],
            [
                "Validation",
                "ast + ruff + black + bandit",
                "Parsed, never imported — importing runs it",
            ],
            ["Database", "SQLAlchemy 2.0", "One connector shape across three SQL dialects"],
            ["Sandbox", "subprocess + rlimits", "Timeout, memory cap, allowlisted environment"],
            [
                "Front end",
                "React 18 + TypeScript + Vite",
                "Built to static files, served by the same process",
            ],
            ["Tests", "pytest — 280 of them", "Including deliberate fault injection"],
        ],
        note="No orchestration framework. The pipeline is six modules and a function call chain.",
    )

    # --- 6-11. The working product -------------------------------------------

    image_slide(
        prs,
        "The product · 1 of 6",
        "Describe the source in plain English",
        shots / "01-welcome.png",
        "No forms, no connection wizard. Three worked examples are built in for a cold start.",
    )
    image_slide(
        prs,
        "The product · 2 of 6",
        "When the request is incomplete, the agent asks",
        shots / "02-clarification.png",
        "Host was never stated, so it is not invented. Confidence drops to 55%, the missing field "
        "is named, and the question carries an example.",
    )
    image_slide(
        prs,
        "The product · 3 of 6",
        "Extraction, generation, and the standards gate",
        shots / "04-standards.png",
        "Everything the agent understood is shown for review — including what it assumed. "
        "13 of 13 checks pass, so the artifact is accepted.",
    )
    image_slide(
        prs,
        "The product · 4 of 6",
        "The generated connector",
        shots / "05-code.png",
        "Provenance above the code: template version, spec checksum, code checksum, repair count. "
        "Written by a template, never by the model.",
    )
    image_slide(
        prs,
        "The product · 5 of 6",
        "Proven against a real database, not just compiled",
        shots / "06-connection.png",
        "Executed in a sandboxed subprocess against live PostgreSQL. Three tables discovered with "
        "their primary keys. Credentials are sent once and never stored, logged, or written to the "
        "artifact.",
    )
    image_slide(
        prs,
        "The product · 6 of 6",
        "The handover artifact",
        shots / "07-artifact.png",
        "Connector, README, requirements, explanation and a manifest, packaged as a versioned zip. "
        "The model reads intent; the template writes the code; the checker decides what ships.",
    )

    # --- 12. Longevity -------------------------------------------------------

    bullets_slide(
        prs,
        "Longevity",
        "Six months later, when someone asks why",
        [
            (
                "Every artifact carries four independent identities",
                "The template version that shaped it, a checksum of the request it came from, a "
                "checksum of the exact bytes emitted, and a version whose patch number is the "
                "repair count — so 1.0.2 reads as “template 1.0, needed two repair passes”, and a "
                "connector that struggled is visible at a glance.",
            ),
            (
                "Packaging is byte-reproducible",
                "The bundle is written with a fixed timestamp rather than the clock, so building "
                "the same artifact twice produces identical bytes. “Is this the file we shipped?” "
                "becomes a comparison instead of a memory.",
            ),
            (
                "The parts that change most often are data, not code",
                "Source types are generated from a JSON file at import — one entry brings the "
                "dialect profile, the clarifying-question options and the registry key with it. "
                "Adding a fourteenth standards check is one function and one list entry.",
            ),
            (
                "Where that stops being true, it is written down",
                "Oracle is a config change; Snowflake needs a new authentication method as well, "
                "because key-pair and SSO do not fit the username/password shape. Saying which is "
                "which is more useful than claiming both are easy.",
            ),
        ],
        note="None of this depends on someone remembering to fill it in — provenance is recorded "
        "by the pipeline that produces the artifact.",
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--shots", type=Path, default=root / "docs" / "screenshots")
    ap.add_argument("--out", type=Path, default=root / "docs" / "Data-Source-Onboarding-Agent.pptx")
    args = ap.parse_args()

    missing = [
        name
        for name in (
            "01-welcome.png",
            "02-clarification.png",
            "04-standards.png",
            "05-code.png",
            "06-connection.png",
            "07-artifact.png",
            "08-manifest.png",
        )
        if not (args.shots / name).is_file()
    ]
    if missing:
        raise SystemExit(f"Missing screenshots in {args.shots}: {', '.join(missing)}")

    path = build(args.shots, args.out)
    print(f"Wrote {path} ({path.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    main()
